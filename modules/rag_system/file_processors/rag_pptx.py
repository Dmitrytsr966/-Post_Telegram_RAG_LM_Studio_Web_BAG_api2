import logging
import re
from typing import Optional, Dict, Any, List

class PPTXCleaner:
    """
    Очиститель текста слайдов PowerPoint:
      - удаляет html-теги
      - удаляет спецсимволы и невидимые символы
      - убирает дублирующиеся пустые строки
      - фильтрует мусорные строки (например, только спецсимволы или пустые)
    """

    def __init__(
        self,
        remove_html: bool = True,
        remove_zero_width: bool = True,
        remove_nbsp: bool = True,
        min_line_length: int = 1,
        junk_patterns: Optional[List[str]] = None,
        logger: Optional[logging.Logger] = None,
    ):
        self.remove_html = remove_html
        self.remove_zero_width = remove_zero_width
        self.remove_nbsp = remove_nbsp
        self.min_line_length = min_line_length
        self.junk_patterns = junk_patterns or []
        self.logger = logger or logging.getLogger(self.__class__.__name__)
        self.ZERO_WIDTH_RE = r'[\u200B-\u200D\uFEFF]'

    def clean_line(self, line: str) -> str:
        s = line
        if self.remove_nbsp:
            s = s.replace('\xa0', ' ').replace('&nbsp;', ' ')
        if self.remove_zero_width:
            s = re.sub(self.ZERO_WIDTH_RE, '', s)
        if self.remove_html:
            s = re.sub(r'<[^>]+>', '', s)
        s = s.strip()
        return s

    def is_junk(self, line: str) -> bool:
        s = line.strip()
        if not s or len(s) < self.min_line_length:
            return True
        if re.fullmatch(r'[-=*_~\s]+', s):
            return True
        for pat in self.junk_patterns:
            if re.fullmatch(pat, s):
                return True
        return False

    def clean_text(self, lines: List[str]) -> List[str]:
        cleaned = []
        for line in lines:
            s = self.clean_line(line)
            if not self.is_junk(s):
                cleaned.append(s)
        return cleaned

class PPTXFileProcessor:
    """
    Класс для обработки pptx-файлов:
      - Читает слайды
      - Очищает каждую строку с помощью PPTXCleaner
      - Собирает результат, пишет логи, возвращает итоговый dict
    """

    def __init__(self, cleaner: Optional[PPTXCleaner] = None, logger: Optional[logging.Logger] = None):
        self.logger = logger or logging.getLogger(self.__class__.__name__)
        self.cleaner = cleaner or PPTXCleaner(logger=self.logger)

    def extract_text(self, file_path: str, **kwargs) -> Dict[str, Any]:
        meta = {
            "file_path": file_path,
            "file_type": "pptx",
            "parser": "rag_pptx",
            "lines": 0,
            "chars": 0,
            "cleaned": False
        }
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slides_text.append(shape.text)
            # Clean and filter lines
            cleaned_lines = self.cleaner.clean_text(slides_text)
            cleaned_text = "\n".join(cleaned_lines)
            meta["lines"] = len(cleaned_lines)
            meta["chars"] = len(cleaned_text)
            meta["cleaned"] = True
            self.logger.info(
                f"PPTX processed: {file_path}, slides: {len(prs.slides)}, lines: {meta['lines']}, chars: {meta['chars']}"
            )
            return {
                "text": cleaned_text,
                "success": True,
                "error": None,
                "cleaned": True,
                "meta": meta
            }
        except Exception as e:
            meta["error"] = str(e)
            self.logger.warning(f"PPTX extraction failed for {file_path}: {e}")
            return {
                "text": "",
                "success": False,
                "error": str(e),
                "cleaned": False,
                "meta": meta
            }

# Адаптер для совместимости с пайплайном
_default_processor = None

def extract_text(file_path: str, cleaner: PPTXCleaner = None, **kwargs) -> Dict[str, Any]:
    global _default_processor
    if _default_processor is None:
        _default_processor = PPTXFileProcessor(cleaner=cleaner)
    return _default_processor.extract_text(file_path, **kwargs)
